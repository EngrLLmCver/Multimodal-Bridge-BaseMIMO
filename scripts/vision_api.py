#!/usr/bin/env python3
"""
multimodal-bridge: 调用 mimo 2.5 多模态 API 提取图片/文档内容为 Markdown。

用法:
  python vision_api.py --file <文件路径> --output <输出路径>
  python vision_api.py --file picture/screenshot.png --output document/screenshot.md
"""

import argparse
import base64
import json
import mimetypes
import os
import sys
from pathlib import Path
from typing import Tuple

# Windows 控制台 UTF-8 输出
if sys.platform == "win32":
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")


# ──────────────────────────────────────────────
# 配置加载
# ──────────────────────────────────────────────

DEFAULT_CONFIG_PATH = r"C:\Users\Administrator\.claude\settings.json"


def load_config(config_path: str = None) -> dict:
    """从 settings.json 读取 API 配置。"""
    path = Path(config_path or DEFAULT_CONFIG_PATH)
    if not path.exists():
        print(f"[ERROR] 配置文件不存在: {path}", file=sys.stderr)
        sys.exit(1)

    with open(path, "r", encoding="utf-8") as f:
        settings = json.load(f)

    env = settings.get("env", {})
    base_url = env.get("ANTHROPIC_BASE_URL", "")
    api_key = env.get("ANTHROPIC_AUTH_TOKEN", "")

    if not base_url or not api_key:
        print("[ERROR] settings.json 中缺少 ANTHROPIC_BASE_URL 或 ANTHROPIC_AUTH_TOKEN", file=sys.stderr)
        sys.exit(1)

    # 确保 base_url 以 /v1/messages 为终点的拼接基础
    base_url = base_url.rstrip("/")

    return {"base_url": base_url, "api_key": api_key}


# ──────────────────────────────────────────────
# 文件类型判断
# ──────────────────────────────────────────────

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".ico", ".tiff", ".tif"}
# SVG 是 XML 文本，当作文档处理而非图片
TEXT_AS_IMAGE_EXTENSIONS = {".svg"}
DOC_EXTENSIONS = {".pdf", ".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt"}


def get_file_type(filepath: str) -> str:
    """返回 'image'、'document' 或 'unknown'。"""
    ext = Path(filepath).suffix.lower()
    if ext in IMAGE_EXTENSIONS:
        return "image"
    if ext in DOC_EXTENSIONS or ext in TEXT_AS_IMAGE_EXTENSIONS:
        return "document"
    return "unknown"


# ──────────────────────────────────────────────
# 图片处理
# ──────────────────────────────────────────────

def encode_image_base64(filepath: str) -> Tuple[str, str]:
    """读取图片并 base64 编码，返回 (base64_data, media_type)。"""
    mime_type, _ = mimetypes.guess_type(filepath)
    if not mime_type:
        mime_type = "image/png"

    file_size = os.path.getsize(filepath)
    # base64 编码后体积增约 33%，Anthropic API 限制约 20MB
    if file_size > 15 * 1024 * 1024:
        print(f"[WARN] 文件较大 ({file_size / 1024 / 1024:.1f}MB)，API 可能超时或拒绝", file=sys.stderr)

    with open(filepath, "rb") as f:
        data = base64.standard_b64encode(f.read()).decode("utf-8")

    return data, mime_type


# ──────────────────────────────────────────────
# 文档处理
# ──────────────────────────────────────────────

def extract_pdf(filepath: str) -> str:
    """用 pdfplumber 提取 PDF 文本。"""
    try:
        import pdfplumber
    except ImportError:
        print("[ERROR] 缺少 pdfplumber，请运行: D:/anaconda3/python.exe -m pip install pdfplumber", file=sys.stderr)
        sys.exit(1)

    texts = []
    with pdfplumber.open(filepath) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text() or ""
            tables = page.extract_tables() or []

            parts = []
            if text.strip():
                parts.append(text)
            for table in tables:
                # 转为 Markdown 表格
                if table and len(table) > 0:
                    header = table[0]
                    md_lines = ["| " + " | ".join(str(c or "") for c in header) + " |"]
                    md_lines.append("| " + " | ".join("---" for _ in header) + " |")
                    for row in table[1:]:
                        md_lines.append("| " + " | ".join(str(c or "") for c in row) + " |")
                    parts.append("\n".join(md_lines))

            if parts:
                texts.append(f"<!-- Page {i+1} -->\n\n" + "\n\n".join(parts))

    return "\n\n---\n\n".join(texts) if texts else "[未能提取到文本内容]"


def extract_docx(filepath: str) -> str:
    """用 python-docx 提取 DOCX 文本。"""
    try:
        from docx import Document
    except ImportError:
        # 回退到 pandoc
        import subprocess
        result = subprocess.run(
            ["pandoc", filepath, "-t", "markdown", "--wrap=none"],
            capture_output=True, text=True, encoding="utf-8"
        )
        if result.returncode == 0:
            return result.stdout
        print("[ERROR] 缺少 python-docx 和 pandoc，无法处理 DOCX", file=sys.stderr)
        sys.exit(1)

    doc = Document(filepath)
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        rows = []
        for row in table.rows:
            rows.append("| " + " | ".join(cell.text for cell in row.cells) + " |")
        if rows:
            # 添加表头分隔线
            rows.insert(1, "| " + " | ".join("---" for _ in table.rows[0].cells) + " |")
            parts.append("\n".join(rows))

    return "\n\n".join(parts) if parts else "[未能提取到文本内容]"


def extract_xlsx(filepath: str) -> str:
    """用 pandas 提取 Excel 内容。"""
    try:
        import pandas as pd
    except ImportError:
        print("[ERROR] 缺少 pandas，请运行: D:/anaconda3/python.exe -m pip install pandas openpyxl", file=sys.stderr)
        sys.exit(1)

    sheets = pd.read_excel(filepath, sheet_name=None, dtype=str)
    parts = []
    for name, df in sheets.items():
        parts.append(f"## Sheet: {name}\n\n{df.fillna('').to_markdown(index=False)}")

    return "\n\n".join(parts) if parts else "[未能提取到内容]"


def extract_pptx(filepath: str) -> str:
    """用 python-pptx 提取 PPTX 文本。"""
    try:
        from pptx import Presentation
    except ImportError:
        print("[ERROR] 缺少 python-pptx，请运行: D:/anaconda3/python.exe -m pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(filepath)
    parts = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        slide_text.append(para.text)
        if slide_text:
            parts.append(f"## Slide {i+1}\n\n" + "\n".join(slide_text))

    return "\n\n".join(parts) if parts else "[未能提取到文本内容]"


def extract_svg(filepath: str) -> str:
    """SVG 是 XML 文本，直接读取。"""
    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        content = f.read()
    return f"```svg\n{content}\n```"


def extract_document_text(filepath: str) -> str:
    """根据文件类型分派提取。"""
    ext = Path(filepath).suffix.lower()
    extractors = {
        ".pdf": extract_pdf,
        ".docx": extract_docx,
        ".doc": extract_docx,
        ".xlsx": extract_xlsx,
        ".xls": extract_xlsx,
        ".pptx": extract_pptx,
        ".ppt": extract_pptx,
        ".svg": extract_svg,
    }
    extractor = extractors.get(ext)
    if not extractor:
        print(f"[ERROR] 不支持的文档格式: {ext}", file=sys.stderr)
        sys.exit(1)
    return extractor(filepath)


# ──────────────────────────────────────────────
# mimo 2.5 API 调用
# ──────────────────────────────────────────────

def _post_to_mimo(base_url: str, api_key: str, payload: dict) -> str:
    """通用 API 调用，自动尝试两种 endpoint 路径，返回文本结果。"""
    import requests

    endpoints = [
        f"{base_url}/v1/messages",
        f"{base_url.rstrip('/anthropic')}/v1/messages",
    ]

    headers = {
        "x-api-key": api_key,
        "anthropic-version": "2023-06-01",
        "content-type": "application/json",
    }

    last_error = None
    for endpoint in endpoints:
        try:
            resp = requests.post(endpoint, json=payload, headers=headers, timeout=120)
            if resp.status_code == 200:
                data = resp.json()
                content = data.get("content", [])
                text_parts = [block["text"] for block in content if block.get("type") == "text"]
                return "\n".join(text_parts) if text_parts else "[API 返回为空]"
            elif resp.status_code == 404:
                continue
            else:
                last_error = f"HTTP {resp.status_code}: {resp.text[:500]}"
        except requests.exceptions.RequestException as e:
            last_error = str(e)

    print(f"[ERROR] API 调用失败: {last_error}", file=sys.stderr)
    sys.exit(1)


def call_mimo_vision(base_url: str, api_key: str, image_b64: str, media_type: str, prompt: str) -> str:
    """调用 mimo 2.5 vision API 识别图片。"""
    payload = {
        "model": "mimo-v2.5",
        "max_tokens": 4096,
        "messages": [{
            "role": "user",
            "content": [
                {
                    "type": "image",
                    "source": {
                        "type": "base64",
                        "media_type": media_type,
                        "data": image_b64,
                    },
                },
                {
                    "type": "text",
                    "text": prompt,
                },
            ],
        }],
    }
    return _post_to_mimo(base_url, api_key, payload)


def call_mimo_text(base_url: str, api_key: str, text_content: str, prompt: str) -> str:
    """调用 mimo 2.5 API 整理文档文本（纯文本模式）。"""
    payload = {
        "model": "mimo-v2.5",
        "max_tokens": 4096,
        "messages": [{
            "role": "user",
            "content": f"{prompt}\n\n以下是原始提取内容：\n\n{text_content}",
        }],
    }
    return _post_to_mimo(base_url, api_key, payload)


# ──────────────────────────────────────────────
# 主流程
# ──────────────────────────────────────────────

IMAGE_PROMPT = """请详细描述这张图片的所有内容，输出为结构化 Markdown：

1. **内容类型**：判断图片类型（照片/截图/图表/文档扫描件/手绘/UI设计/流程图/其他）
2. **视觉描述**：整体布局、颜色、构图
3. **文字内容**：完整转录图片中所有文字，保持原始排版和层级
4. **数据/表格**：如包含表格或图表数据，转为 Markdown 表格
5. **关键信息**：列出要点

用中文输出，保留原文中的英文术语。对于无法识别的内容标注 [无法识别]。"""

DOC_PROMPT = """以下是用 Python 库从文档中提取的原始文本。请将其整理为结构化 Markdown：

1. 修正明显的格式问题和乱码
2. 保留原始层级结构（标题、段落、列表）
3. 表格数据确保是正确的 Markdown 表格格式
4. 如果内容有分页，保留分页标记
5. 在文件开头添加一段简要摘要

用中文输出摘要部分，正文保留原始语言。"""


def main():
    parser = argparse.ArgumentParser(description="multimodal-bridge: 多模态内容提取为 Markdown")
    parser.add_argument("--file", required=True, help="输入文件路径（图片或文档）")
    parser.add_argument("--output", required=True, help="Markdown 输出路径")
    parser.add_argument("--config", default=None, help="配置文件路径（默认: settings.json）")
    parser.add_argument("--prompt", default=None, help="自定义提示词（覆盖默认）")
    args = parser.parse_args()

    # 检查输入文件
    input_path = Path(args.file)
    if not input_path.exists():
        print(f"[ERROR] 文件不存在: {input_path}", file=sys.stderr)
        sys.exit(1)

    # 确保输出目录存在
    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    # 加载配置
    config = load_config(args.config)

    # 判断文件类型
    file_type = get_file_type(str(input_path))
    print(f"[INFO] 文件类型: {file_type}")
    print(f"[INFO] 输入: {input_path}")
    print(f"[INFO] 输出: {output_path}")

    if file_type == "image":
        # 图片：base64 编码 → mimo 2.5 vision API
        print("[INFO] 正在编码图片...")
        image_b64, media_type = encode_image_base64(str(input_path))

        print("[INFO] 正在调用 mimo 2.5 vision API...")
        prompt = args.prompt or IMAGE_PROMPT
        result = call_mimo_vision(config["base_url"], config["api_key"], image_b64, media_type, prompt)

    elif file_type == "document":
        # 文档：先提取文本 → mimo 2.5 整理
        print("[INFO] 正在提取文档文本...")
        raw_text = extract_document_text(str(input_path))

        print("[INFO] 正在调用 mimo 2.5 整理内容...")
        prompt = args.prompt or DOC_PROMPT
        result = call_mimo_text(config["base_url"], config["api_key"], raw_text, prompt)

    else:
        print(f"[ERROR] 不支持的文件格式: {input_path.suffix}", file=sys.stderr)
        sys.exit(1)

    # 写入输出文件
    output_path.write_text(result, encoding="utf-8")
    print(f"[OK] 已保存至 {output_path}")


if __name__ == "__main__":
    main()
